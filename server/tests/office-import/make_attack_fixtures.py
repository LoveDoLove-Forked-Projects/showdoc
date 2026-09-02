#!/usr/bin/env python3
"""
生成 Office 导入安全攻击样本（供 http_security_test.php 使用）。纯 python，流式写 zip。

输出到 fixtures/（可用 FX_OUT 覆盖）：
  atk_fakeext.docx    HTML+PHP 内容伪造 .docx（魔数校验应拒绝）
  atk_polyglot.docx   PHP 前缀 + 真实 zip 混合（魔数校验应拒绝）
  atk_xss.docx        正文注入 <script>/<img onerror>/javascript:（净化应中和）
  atk_longtitle.docx  H1/H2 超 50 字符（截断应成功写库）
  atk_media.pptx      media 含路径穿越条目 + xss.svg + 伪造 png(php)，且被 slide rels 引用
  atk_bomb.docx       2000 x 512KB 重复条目（解压总量 ~1GB，超上限）
"""
import os
import sys
import zipfile

try:
    import docx
    import pptx
except ImportError:
    print("需要 python-docx / python-pptx（宿主机）", file=sys.stderr)
    sys.exit(1)

out = os.environ.get("FX_OUT") or os.path.join(os.path.dirname(os.path.abspath(__file__)), "fixtures")
os.makedirs(out, exist_ok=True)

# 1. 伪造扩展名
with open(os.path.join(out, "atk_fakeext.docx"), "w") as f:
    f.write('<html><body><?php echo "webshell"; ?></body></html>')
print("atk_fakeext.docx")

# 2. polyglot：先造真 docx 再前置 PHP
doc = docx.Document()
doc.add_paragraph("innocent")
base = os.path.join(out, "_poly_base.docx")
doc.save(base)
with open(base, "rb") as f:
    zbytes = f.read()
with open(os.path.join(out, "atk_polyglot.docx"), "wb") as f:
    f.write(b"<?php system($_GET['c']); ?>\n" + zbytes)
os.unlink(base)
print("atk_polyglot.docx")

# 3. XSS 注入
doc = docx.Document()
doc.add_heading("XSS Attack", level=1)
for line in [
    '<script>alert("doc-xss")</script>',
    '<img src=x onerror=alert(2)>',
    '[click me](javascript:alert(3))',
    '![x](javascript:alert(4))',
    '<iframe src="javascript:alert(5)"></iframe>',
    '<svg onload=alert(6)></svg>',
    '[实体混淆](&#106;avascript:alert(7))',
    '正文正常文字，含 [正常链接](https://showdoc.com.cn) 与行内 `code`。',
    '```',
    '<script>alert("in-code-fence-must-stay")</script>',
    '```',
]:
    doc.add_paragraph(line)
doc.save(os.path.join(out, "atk_xss.docx"))
print("atk_xss.docx")

# 4. 超长标题
doc = docx.Document()
doc.add_heading("标" * 120, level=1)
doc.add_paragraph("内容 A")
doc.add_heading("子" * 100, level=2)
doc.add_paragraph("内容 B")
doc.save(os.path.join(out, "atk_longtitle.docx"))
print("atk_longtitle.docx")

# 5. 恶意 media 的 pptx
prs = pptx.Presentation()
s = prs.slides.add_slide(prs.slide_layouts[5])
s.shapes.title.text = "Evil Media"
p = os.path.join(out, "atk_media.pptx")
prs.save(p)
# 塞入：路径穿越条目、带脚本的 svg、伪 png（实为 PHP），并把 rId901/rId902 挂进 slide1
with zipfile.ZipFile(p, "a") as z:
    z.writestr("../../../var/www/shell.php", "<?php evil(); ?>")
    z.writestr("ppt/media/xss.svg", '<svg xmlns="http://www.w3.org/2000/svg" onload="alert(1)"><script>alert(2)</script></svg>')
    z.writestr("ppt/media/fake.png", '<?php echo "not a png"; ?>')
    rels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
    rels = rels.replace(
        "</Relationships>",
        '<Relationship Id="rId901" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/xss.svg"/>'
        '<Relationship Id="rId902" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/fake.png"/>'
        "</Relationships>",
    )
    slide = z.read("ppt/slides/slide1.xml").decode()
    blip = '<a:blip r:embed="rId901" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/>'
    blip += '<a:blip r:embed="rId902" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/>'
    slide = slide.replace("</p:sld>", blip + "</p:sld>")
import shutil, tempfile
tmp = tempfile.mktemp(suffix=".pptx")
with zipfile.ZipFile(p) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        if item.filename == "ppt/slides/_rels/slide1.xml.rels":
            zout.writestr(item, rels)
        elif item.filename == "ppt/slides/slide1.xml":
            zout.writestr(item, slide)
        else:
            zout.writestr(item, zin.read(item.filename))
shutil.move(tmp, p)
print("atk_media.pptx")

# 6. zip 炸弹（2000 x 512KB 全零，压缩后极小，解压声明总量 ~1GB）
p = os.path.join(out, "atk_bomb.docx")
chunk = "\x00" * (512 * 1024)
with zipfile.ZipFile(p, "w", zipfile.ZIP_DEFLATED) as z:
    z.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types/>')
    z.writestr("word/document.xml", '<?xml version="1.0"?><w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p/></w:body></w:document>')
    for i in range(2000):
        z.writestr(f"word/chunk_{i}.bin", chunk)
print("atk_bomb.docx")
